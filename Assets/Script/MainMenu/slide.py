from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# --- CONFIGURATION & COLORS ---
# Neobrutalist Palette based on your screenshots
COLOR_ACID_GREEN = RGBColor(204, 255, 0)   # The bright green
COLOR_LAVENDER = RGBColor(160, 107, 251)   # The purple/lilac
COLOR_BLACK = RGBColor(0, 0, 0)            # Borders and Shadows
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_BG = RGBColor(248, 255, 230)         # Very pale background yellow/green tint

# Create Presentation
prs = Presentation()
prs.slide_width = Inches(13.333) # Wide format
prs.slide_height = Inches(7.5)

# --- HELPER FUNCTIONS ---

def set_background(slide):
    """Sets the slide background to the pale wavy-ish color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_neobrutalist_box(slide, x, y, w, h, text, bg_color, font_size=24, bold=True, alignment=PP_ALIGN.CENTER):
    """
    Creates a box with a thick black border and a hard black drop shadow (offset).
    """
    # 1. Shadow (Black box offset)
    shadow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.1), y + Inches(0.1), w, h)
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = COLOR_BLACK
    shadow.line.color.rgb = COLOR_BLACK
    
    # 2. Main Box
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    
    # Thick border
    line = box.line
    line.color.rgb = COLOR_BLACK
    line.width = Pt(4.5) 

    # Text
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = alignment
    p.font.family = "Arial Black" if bold else "Arial"
    p.font.size = Pt(font_size)
    p.font.color.rgb = COLOR_BLACK
    p.font.bold = bold
    
    return box

def add_pixel_text(slide, x, y, w, h, text, size=18):
    """Adds text that mimics the pixelated font style."""
    textbox = slide.shapes.add_textbox(x, y, w, h)
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.family = "Courier New" # Monospace to simulate pixel/code look
    p.font.bold = True
    p.font.size = Pt(size)
    p.font.color.rgb = COLOR_BLACK
    return textbox

# --- SLIDE 1: TITLE ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)

# Title Elements
add_neobrutalist_box(slide, Inches(3.5), Inches(1.5), Inches(6.3), Inches(1.5), "anchOUR", COLOR_ACID_GREEN, font_size=80)
add_neobrutalist_box(slide, Inches(2), Inches(3.5), Inches(9.33), Inches(1.0), "Don't let your Community Hit Rock Bottom", COLOR_LAVENDER, font_size=32)

# Decorative "OUI" badge
add_neobrutalist_box(slide, Inches(10.5), Inches(0.5), Inches(1.5), Inches(0.8), "OUI", COLOR_LAVENDER, font_size=24)

# Team Names at bottom
add_pixel_text(slide, Inches(1), Inches(6.5), Inches(11), Inches(1), 
               "TEAM: Yutong Wang | Matthew Gorman | Stefaniia Kornilova | Elizabeth")


# --- SLIDE 2: INSPIRATION (The Problem) ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)

# Header
add_neobrutalist_box(slide, Inches(0.5), Inches(0.5), Inches(4), Inches(1), "The Problem", COLOR_ACID_GREEN, font_size=36)

# Content - Left Side
add_pixel_text(slide, Inches(0.8), Inches(2), Inches(6), Inches(4), 
               "Community drift is silent.\n\n"
               "Interest slowly fades without anyone noticing.\n\n"
               "By the time leaders realize engagement has collapsed,\n"
               "it is often too late to recover.")

# Visual - Right Side (Symbolic Sinking Anchor)
# Using a box to represent the "Anchor" visual area
box = add_neobrutalist_box(slide, Inches(7.5), Inches(2.5), Inches(5), Inches(3), "ENGAGEMENT DECAY", COLOR_WHITE, font_size=24)
# Add a downward trend line manually
line = slide.shapes.add_connector(MSO_SHAPE.STRAIGHT_CONNECTOR_1, Inches(8), Inches(3), Inches(12), Inches(5))
line.line.color.rgb = COLOR_LAVENDER
line.line.width = Pt(8)


# --- SLIDE 3: THE SOLUTION ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)

# Header
add_neobrutalist_box(slide, Inches(0.5), Inches(0.5), Inches(4), Inches(1), "What It Does", COLOR_LAVENDER, font_size=36)

# Three Cards feature layout
# Card 1: Analyze
add_neobrutalist_box(slide, Inches(1), Inches(2.5), Inches(3.5), Inches(3.5), "Analyze\nDecay", COLOR_WHITE, font_size=28)
add_pixel_text(slide, Inches(1.2), Inches(4.5), Inches(3), Inches(2), "Visualizes attendance\nand follower drop-offs.")

# Card 2: Explain
add_neobrutalist_box(slide, Inches(4.9), Inches(2.5), Inches(3.5), Inches(3.5), "Explain\nWhy", COLOR_ACID_GREEN, font_size=28)
add_pixel_text(slide, Inches(5.1), Inches(4.5), Inches(3), Inches(2), "Detects declining\nparticipation signals.")

# Card 3: Stabilize
add_neobrutalist_box(slide, Inches(8.8), Inches(2.5), Inches(3.5), Inches(3.5), "Secure\nFuture", COLOR_WHITE, font_size=28)
add_pixel_text(slide, Inches(9.0), Inches(4.5), Inches(3), Inches(2), "Data-backed insights\nto re-engage members.")


# --- SLIDE 4: HOW WE BUILT IT (Tech Stack) ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)

# Header
add_neobrutalist_box(slide, Inches(0.5), Inches(0.5), Inches(5), Inches(1), "Tech Stack", COLOR_ACID_GREEN, font_size=36)

# Architecture Flow
# 1. Data Collection
add_neobrutalist_box(slide, Inches(1), Inches(2.5), Inches(2.5), Inches(1.5), "Python +\nApify", COLOR_WHITE, font_size=20)
add_pixel_text(slide, Inches(1), Inches(4.1), Inches(2.5), Inches(1), "Data Collection", 12)

# Arrow
arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.6), Inches(3), Inches(0.8), Inches(0.5))
arrow.fill.solid()
arrow.fill.fore_color.rgb = COLOR_BLACK

# 2. Storage & ML
add_neobrutalist_box(slide, Inches(4.5), Inches(2), Inches(4.3), Inches(2.5), "SNOWFLAKE", COLOR_LAVENDER, font_size=30)
add_pixel_text(slide, Inches(4.6), Inches(3.5), Inches(4), Inches(1), 
               "- Cortex AI + Gemini (Insights)\n- Arctic-embed (Vectors)\n- Scikit-learn (Clustering)")

# Arrow
arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.9), Inches(3), Inches(0.8), Inches(0.5))
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = COLOR_BLACK

# 3. Frontend
add_neobrutalist_box(slide, Inches(9.8), Inches(2.5), Inches(2.5), Inches(1.5), "Streamlit", COLOR_WHITE, font_size=20)
add_pixel_text(slide, Inches(9.8), Inches(4.1), Inches(2.5), Inches(1), "Interactive Dashboard", 12)

# Footer Note
add_pixel_text(slide, Inches(4), Inches(6.5), Inches(6), Inches(1), "Design generated by Google AI Studio", 14)


# --- SLIDE 5: DEMO PLACEHOLDER ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)

# Big Center Button for Demo
add_neobrutalist_box(slide, Inches(3.66), Inches(2.5), Inches(6), Inches(2.5), "LIVE DEMO", COLOR_ACID_GREEN, font_size=60)
add_pixel_text(slide, Inches(4), Inches(5.2), Inches(5), Inches(1), "Please switch screens...", 18)


# --- SLIDE 6: WHAT'S NEXT & TEAM ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)

# Header
add_neobrutalist_box(slide, Inches(0.5), Inches(0.5), Inches(4.5), Inches(1), "What's Next?", COLOR_LAVENDER, font_size=36)

# Future Content
box_future = add_neobrutalist_box(slide, Inches(1), Inches(2), Inches(11.3), Inches(2), "", COLOR_WHITE, font_size=1)
tf = box_future.text_frame
tf.clear()
p = tf.paragraphs[0]
p.text = "Full integration with Meta API"
p.font.size = Pt(28)
p.font.bold = True
p.font.family = "Arial Black"
p.font.color.rgb = COLOR_BLACK

p2 = tf.add_paragraph()
p2.text = "Real-time Instagram engagement data & stronger predictive signals."
p2.font.size = Pt(20)
p2.font.family = "Courier New"
p2.font.color.rgb = COLOR_BLACK

# Accomplishments / Proud Of
add_pixel_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(2),
               "PROUD OF:\n"
               "- End-to-End Analytics Pipeline inside Snowflake\n"
               "- Turning abstract 'Health' into measurable signals\n"
               "- Neobrutalist design implementation")

# Team
add_neobrutalist_box(slide, Inches(9), Inches(6), Inches(3.5), Inches(1), "Thank You!", COLOR_ACID_GREEN, font_size=24)

# Save
prs.save('anchOUR_Presentation.pptx')
print("Presentation saved as anchOUR_Presentation.pptx")